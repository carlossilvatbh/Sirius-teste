"""
Organogram Printing Views
Professional printing and report generation for organograms
"""

from django.shortcuts import render, get_object_or_404
from django.http import JsonResponse, HttpResponse
from django.contrib.admin.views.decorators import staff_member_required
from django.template.loader import render_to_string
from django.views.decorators.csrf import csrf_exempt
from django.conf import settings
import json
import os
from datetime import datetime
import base64
from io import BytesIO

# Import for PDF generation
try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4, legal
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

from .models import Entity, Structure, EntityOwnership
from parties.models import Party


@staff_member_required
def organogram_printing_view(request):
    """
    Main organogram printing interface
    """
    structures = Structure.objects.all().order_by('name')
    parties = Party.objects.all().order_by('name')
    
    context = {
        'structures': structures,
        'parties': parties,
        'print_templates': get_print_templates(),
        'export_formats': get_export_formats()
    }
    
    return render(request, 'admin/corporate/organogram_printing.html', context)


@staff_member_required
@csrf_exempt
def generate_organogram_report(request):
    """
    Generate organogram report in various formats
    """
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Method not allowed'})
    
    try:
        data = json.loads(request.body)
        
        structure_id = data.get('structure_id')
        party_id = data.get('party_id')
        report_type = data.get('report_type', 'structure_overview')
        format_type = data.get('format', 'pdf')
        template = data.get('template', 'professional')
        options = data.get('options', {})
        
        # Generate report based on parameters
        if structure_id:
            structure = get_object_or_404(Structure, id=structure_id)
            report_data = generate_structure_report_data(structure, options)
            filename = f"organogram_{structure.name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        elif party_id:
            party = get_object_or_404(Party, id=party_id)
            report_data = generate_party_report_data(party, options)
            filename = f"party_organogram_{party.name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        else:
            return JsonResponse({'success': False, 'error': 'Structure or Party ID required'})
        
        # Generate report in requested format
        if format_type == 'pdf':
            file_path = generate_pdf_report(report_data, template, filename, options)
        elif format_type == 'html':
            file_path = generate_html_report(report_data, template, filename, options)
        elif format_type == 'excel':
            file_path = generate_excel_report(report_data, filename, options)
        elif format_type == 'powerpoint':
            file_path = generate_powerpoint_report(report_data, template, filename, options)
        else:
            return JsonResponse({'success': False, 'error': 'Unsupported format'})
        
        return JsonResponse({
            'success': True,
            'message': 'Report generated successfully',
            'download_url': f'/admin/corporate/download-report/{os.path.basename(file_path)}/',
            'filename': os.path.basename(file_path)
        })
        
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': str(e)
        })


@staff_member_required
def download_report(request, filename):
    """
    Download generated report
    """
    try:
        file_path = os.path.join(settings.MEDIA_ROOT, 'reports', filename)
        
        if not os.path.exists(file_path):
            return JsonResponse({'error': 'File not found'}, status=404)
        
        with open(file_path, 'rb') as f:
            response = HttpResponse(f.read())
            
        # Set appropriate content type
        if filename.endswith('.pdf'):
            response['Content-Type'] = 'application/pdf'
        elif filename.endswith('.html'):
            response['Content-Type'] = 'text/html'
        elif filename.endswith('.xlsx'):
            response['Content-Type'] = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        elif filename.endswith('.pptx'):
            response['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        
        response['Content-Disposition'] = f'attachment; filename="{filename}"'
        return response
        
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@staff_member_required
def preview_organogram(request):
    """
    Generate preview of organogram for printing
    """
    structure_id = request.GET.get('structure_id')
    party_id = request.GET.get('party_id')
    template = request.GET.get('template', 'professional')
    
    try:
        if structure_id:
            structure = get_object_or_404(Structure, id=structure_id)
            preview_data = generate_structure_preview_data(structure)
            title = f"Structure: {structure.name}"
        elif party_id:
            party = get_object_or_404(Party, id=party_id)
            preview_data = generate_party_preview_data(party)
            title = f"Party: {party.name}"
        else:
            return JsonResponse({'success': False, 'error': 'Structure or Party ID required'})
        
        context = {
            'title': title,
            'preview_data': preview_data,
            'template': template,
            'generated_at': datetime.now()
        }
        
        return render(request, 'admin/corporate/organogram_preview.html', context)
        
    except Exception as e:
        return JsonResponse({'success': False, 'error': str(e)})


def generate_structure_report_data(structure, options):
    """
    Generate comprehensive data for structure report
    """
    # Get all entities in the structure
    ownerships = EntityOwnership.objects.filter(structure=structure).select_related(
        'owner_entity', 'owner_ubo', 'owned_entity'
    )
    
    # Build hierarchy
    hierarchy = build_ownership_hierarchy(ownerships)
    
    # Calculate statistics
    stats = {
        'total_entities': Entity.objects.filter(
            id__in=ownerships.values_list('owned_entity_id', flat=True)
        ).count(),
        'total_parties': Party.objects.filter(
            id__in=ownerships.values_list('owner_ubo_id', flat=True)
        ).count(),
        'total_ownerships': ownerships.count(),
        'avg_ownership': ownerships.aggregate(avg=models.Avg('ownership_percentage'))['avg'] or 0,
        'max_ownership': ownerships.aggregate(max=models.Max('ownership_percentage'))['max'] or 0
    }
    
    # Get jurisdictions
    entities = Entity.objects.filter(
        id__in=ownerships.values_list('owned_entity_id', flat=True)
    )
    jurisdictions = list(set(
        getattr(entity, 'jurisdiction', 'Unknown')
        for entity in entities
    ))
    
    return {
        'structure': {
            'id': structure.id,
            'name': structure.name,
            'structure_type': getattr(structure, 'structure_type', 'Unknown'),
            'created_date': getattr(structure, 'created_date', datetime.now()),
            'status': 'Active'
        },
        'hierarchy': hierarchy,
        'statistics': stats,
        'jurisdictions': jurisdictions,
        'entities': [
            {
                'id': entity.id,
                'name': entity.name,
                'entity_type': entity.entity_type,
                'jurisdiction': getattr(entity, 'jurisdiction', 'Unknown')
            }
            for entity in entities
        ],
        'ownerships': [
            {
                'owner_name': ownership.owner_ubo.name if ownership.owner_ubo else ownership.owner_entity.name,
                'owner_type': 'Party' if ownership.owner_ubo else 'Entity',
                'owned_entity': ownership.owned_entity.name,
                'ownership_percentage': ownership.ownership_percentage,
                'notes': getattr(ownership, 'notes', '')
            }
            for ownership in ownerships
        ]
    }


def generate_party_report_data(party, options):
    """
    Generate comprehensive data for party report
    """
    # Get all ownerships where this party is the owner
    ownerships = EntityOwnership.objects.filter(owner_ubo=party).select_related(
        'owned_entity', 'structure'
    )
    
    # Get structures
    structures = Structure.objects.filter(
        id__in=ownerships.values_list('structure_id', flat=True)
    ).distinct()
    
    # Calculate statistics
    stats = {
        'total_structures': structures.count(),
        'total_entities': ownerships.values('owned_entity').distinct().count(),
        'total_ownership': ownerships.aggregate(total=models.Sum('ownership_percentage'))['total'] or 0,
        'avg_ownership': ownerships.aggregate(avg=models.Avg('ownership_percentage'))['avg'] or 0,
        'max_ownership': ownerships.aggregate(max=models.Max('ownership_percentage'))['max'] or 0
    }
    
    return {
        'party': {
            'id': party.id,
            'name': party.name,
            'party_type': getattr(party, 'party_type', 'Individual')
        },
        'statistics': stats,
        'structures': [
            {
                'id': structure.id,
                'name': structure.name,
                'structure_type': getattr(structure, 'structure_type', 'Unknown'),
                'party_ownership': ownerships.filter(structure=structure).aggregate(
                    total=models.Sum('ownership_percentage')
                )['total'] or 0
            }
            for structure in structures
        ],
        'ownerships': [
            {
                'structure_name': ownership.structure.name,
                'owned_entity': ownership.owned_entity.name,
                'ownership_percentage': ownership.ownership_percentage,
                'entity_type': ownership.owned_entity.entity_type,
                'jurisdiction': getattr(ownership.owned_entity, 'jurisdiction', 'Unknown')
            }
            for ownership in ownerships
        ]
    }


def build_ownership_hierarchy(ownerships):
    """
    Build hierarchical structure from ownerships
    """
    hierarchy = {}
    
    # Group by owner
    for ownership in ownerships:
        owner_key = ownership.owner_ubo.name if ownership.owner_ubo else ownership.owner_entity.name
        
        if owner_key not in hierarchy:
            hierarchy[owner_key] = {
                'type': 'Party' if ownership.owner_ubo else 'Entity',
                'owned_entities': []
            }
        
        hierarchy[owner_key]['owned_entities'].append({
            'name': ownership.owned_entity.name,
            'ownership_percentage': ownership.ownership_percentage,
            'entity_type': ownership.owned_entity.entity_type
        })
    
    return hierarchy


def generate_pdf_report(report_data, template, filename, options):
    """
    Generate PDF report using ReportLab
    """
    if not REPORTLAB_AVAILABLE:
        raise Exception("ReportLab not available for PDF generation")
    
    # Ensure reports directory exists
    reports_dir = os.path.join(settings.MEDIA_ROOT, 'reports')
    os.makedirs(reports_dir, exist_ok=True)
    
    file_path = os.path.join(reports_dir, f"{filename}.pdf")
    
    # Create PDF document
    doc = SimpleDocTemplate(
        file_path,
        pagesize=A4,
        rightMargin=72,
        leftMargin=72,
        topMargin=72,
        bottomMargin=18
    )
    
    # Build story
    story = []
    styles = getSampleStyleSheet()
    
    # Title
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        spaceAfter=30,
        alignment=TA_CENTER,
        textColor=HexColor('#667eea')
    )
    
    if 'structure' in report_data:
        title = f"Structure Organogram: {report_data['structure']['name']}"
    else:
        title = f"Party Organogram: {report_data['party']['name']}"
    
    story.append(Paragraph(title, title_style))
    story.append(Spacer(1, 20))
    
    # Add content based on template
    if template == 'professional':
        story.extend(build_professional_pdf_content(report_data, styles))
    elif template == 'executive':
        story.extend(build_executive_pdf_content(report_data, styles))
    elif template == 'detailed':
        story.extend(build_detailed_pdf_content(report_data, styles))
    
    # Build PDF
    doc.build(story)
    
    return file_path


def build_professional_pdf_content(report_data, styles):
    """
    Build professional PDF content
    """
    story = []
    
    # Executive Summary
    story.append(Paragraph("Executive Summary", styles['Heading2']))
    story.append(Spacer(1, 12))
    
    if 'structure' in report_data:
        summary_text = f"""
        This report provides a comprehensive overview of the {report_data['structure']['name']} structure.
        The structure contains {report_data['statistics']['total_entities']} entities and 
        {report_data['statistics']['total_parties']} parties with 
        {report_data['statistics']['total_ownerships']} ownership relationships.
        """
    else:
        summary_text = f"""
        This report provides a comprehensive overview of {report_data['party']['name']}'s ownership portfolio.
        The party has interests in {report_data['statistics']['total_structures']} structures and 
        {report_data['statistics']['total_entities']} entities with total ownership of 
        {report_data['statistics']['total_ownership']:.2f}%.
        """
    
    story.append(Paragraph(summary_text, styles['Normal']))
    story.append(Spacer(1, 20))
    
    # Statistics Table
    story.append(Paragraph("Key Statistics", styles['Heading2']))
    story.append(Spacer(1, 12))
    
    stats_data = [['Metric', 'Value']]
    for key, value in report_data['statistics'].items():
        if isinstance(value, float):
            stats_data.append([key.replace('_', ' ').title(), f"{value:.2f}"])
        else:
            stats_data.append([key.replace('_', ' ').title(), str(value)])
    
    stats_table = Table(stats_data)
    stats_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), HexColor('#667eea')),
        ('TEXTCOLOR', (0, 0), (-1, 0), HexColor('#ffffff')),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 14),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), HexColor('#f7fafc')),
        ('GRID', (0, 0), (-1, -1), 1, HexColor('#e2e8f0'))
    ]))
    
    story.append(stats_table)
    story.append(Spacer(1, 20))
    
    # Ownership Details
    if 'ownerships' in report_data:
        story.append(Paragraph("Ownership Details", styles['Heading2']))
        story.append(Spacer(1, 12))
        
        ownership_data = [['Owner', 'Owned Entity', 'Ownership %', 'Type']]
        for ownership in report_data['ownerships'][:20]:  # Limit to first 20
            ownership_data.append([
                ownership['owner_name'],
                ownership['owned_entity'],
                f"{ownership['ownership_percentage']:.2f}%",
                ownership.get('owner_type', 'Entity')
            ])
        
        ownership_table = Table(ownership_data)
        ownership_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), HexColor('#667eea')),
            ('TEXTCOLOR', (0, 0), (-1, 0), HexColor('#ffffff')),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), HexColor('#f7fafc')),
            ('GRID', (0, 0), (-1, -1), 1, HexColor('#e2e8f0'))
        ]))
        
        story.append(ownership_table)
    
    return story


def build_executive_pdf_content(report_data, styles):
    """
    Build executive summary PDF content
    """
    return build_professional_pdf_content(report_data, styles)  # Simplified for now


def build_detailed_pdf_content(report_data, styles):
    """
    Build detailed PDF content
    """
    return build_professional_pdf_content(report_data, styles)  # Simplified for now


def generate_html_report(report_data, template, filename, options):
    """
    Generate HTML report
    """
    reports_dir = os.path.join(settings.MEDIA_ROOT, 'reports')
    os.makedirs(reports_dir, exist_ok=True)
    
    file_path = os.path.join(reports_dir, f"{filename}.html")
    
    context = {
        'report_data': report_data,
        'template': template,
        'generated_at': datetime.now(),
        'options': options
    }
    
    html_content = render_to_string('admin/corporate/organogram_report.html', context)
    
    with open(file_path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    return file_path


def generate_excel_report(report_data, filename, options):
    """
    Generate Excel report
    """
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
    except ImportError:
        raise Exception("openpyxl not available for Excel generation")
    
    reports_dir = os.path.join(settings.MEDIA_ROOT, 'reports')
    os.makedirs(reports_dir, exist_ok=True)
    
    file_path = os.path.join(reports_dir, f"{filename}.xlsx")
    
    # Create workbook
    wb = openpyxl.Workbook()
    ws = wb.active
    
    # Set title
    if 'structure' in report_data:
        ws.title = "Structure Report"
        ws['A1'] = f"Structure: {report_data['structure']['name']}"
    else:
        ws.title = "Party Report"
        ws['A1'] = f"Party: {report_data['party']['name']}"
    
    # Style title
    ws['A1'].font = Font(size=16, bold=True, color='667eea')
    ws['A1'].alignment = Alignment(horizontal='center')
    
    # Add statistics
    row = 3
    ws[f'A{row}'] = "Statistics"
    ws[f'A{row}'].font = Font(size=14, bold=True)
    row += 1
    
    for key, value in report_data['statistics'].items():
        ws[f'A{row}'] = key.replace('_', ' ').title()
        ws[f'B{row}'] = value
        row += 1
    
    # Add ownership data
    row += 2
    ws[f'A{row}'] = "Ownership Details"
    ws[f'A{row}'].font = Font(size=14, bold=True)
    row += 1
    
    # Headers
    headers = ['Owner', 'Owned Entity', 'Ownership %', 'Type']
    for col, header in enumerate(headers, 1):
        cell = ws.cell(row=row, column=col, value=header)
        cell.font = Font(bold=True)
        cell.fill = PatternFill(start_color='667eea', end_color='667eea', fill_type='solid')
    
    row += 1
    
    # Data
    for ownership in report_data.get('ownerships', []):
        ws.cell(row=row, column=1, value=ownership['owner_name'])
        ws.cell(row=row, column=2, value=ownership['owned_entity'])
        ws.cell(row=row, column=3, value=ownership['ownership_percentage'])
        ws.cell(row=row, column=4, value=ownership.get('owner_type', 'Entity'))
        row += 1
    
    # Save workbook
    wb.save(file_path)
    
    return file_path


def generate_powerpoint_report(report_data, template, filename, options):
    """
    Generate PowerPoint report
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        raise Exception("python-pptx not available for PowerPoint generation")
    
    reports_dir = os.path.join(settings.MEDIA_ROOT, 'reports')
    os.makedirs(reports_dir, exist_ok=True)
    
    file_path = os.path.join(reports_dir, f"{filename}.pptx")
    
    # Create presentation
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    if 'structure' in report_data:
        title = f"Structure Organogram"
        subtitle = report_data['structure']['name']
    else:
        title = f"Party Organogram"
        subtitle = report_data['party']['name']
    
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    
    # Statistics slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Key Statistics"
    
    content = slide.placeholders[1].text_frame
    for key, value in report_data['statistics'].items():
        p = content.paragraphs[0] if len(content.paragraphs) == 1 and not content.paragraphs[0].text else content.add_paragraph()
        p.text = f"{key.replace('_', ' ').title()}: {value}"
        p.level = 0
    
    # Save presentation
    prs.save(file_path)
    
    return file_path


def generate_structure_preview_data(structure):
    """
    Generate preview data for structure
    """
    ownerships = EntityOwnership.objects.filter(structure=structure)
    return {
        'type': 'structure',
        'name': structure.name,
        'entity_count': ownerships.values('owned_entity').distinct().count(),
        'ownership_count': ownerships.count()
    }


def generate_party_preview_data(party):
    """
    Generate preview data for party
    """
    ownerships = EntityOwnership.objects.filter(owner_ubo=party)
    return {
        'type': 'party',
        'name': party.name,
        'structure_count': ownerships.values('structure').distinct().count(),
        'ownership_count': ownerships.count()
    }


def get_print_templates():
    """
    Get available print templates
    """
    return [
        {'id': 'professional', 'name': 'Professional', 'description': 'Clean, business-focused layout'},
        {'id': 'executive', 'name': 'Executive Summary', 'description': 'High-level overview for executives'},
        {'id': 'detailed', 'name': 'Detailed Analysis', 'description': 'Comprehensive technical report'},
        {'id': 'compliance', 'name': 'Compliance Report', 'description': 'Regulatory compliance focused'}
    ]


def get_export_formats():
    """
    Get available export formats
    """
    formats = [
        {'id': 'pdf', 'name': 'PDF Document', 'description': 'Portable document format'},
        {'id': 'html', 'name': 'HTML Report', 'description': 'Web-based interactive report'}
    ]
    
    if REPORTLAB_AVAILABLE:
        formats.append({'id': 'excel', 'name': 'Excel Spreadsheet', 'description': 'Microsoft Excel format'})
    
    try:
        import pptx
        formats.append({'id': 'powerpoint', 'name': 'PowerPoint Presentation', 'description': 'Microsoft PowerPoint format'})
    except ImportError:
        pass
    
    return formats

